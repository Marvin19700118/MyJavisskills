from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_emba_presentation():
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI 賦能企業：從技術趨勢到商業落地的實戰指南"
    slide.placeholders[1].text = "EMBA 專屬高階經理人講座\n(預計演講時長：3小時)"

    # Outline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "今日議程 (Agenda)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Module 1: 揭開 AI 神秘面紗 —— 看懂技術本質 (45 mins)"
    tf.add_paragraph().text = "Module 2: 2026 最熱門的 AI 話題與前沿趨勢 (40 mins)"
    tf.add_paragraph().text = "Module 3: 企業實戰 —— AI 帶來的商業價值與成功案例 (50 mins)"
    tf.add_paragraph().text = "Module 4: 高管必修 —— 企業導入 AI 的策略與風險管理 (30 mins)"
    tf.add_paragraph().text = "Q&A 交流時間 (15 mins)"

    # --- Module 1 ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "【Module 1】揭開 AI 神秘面紗"
    slide.placeholders[1].text = "核心目標：拆解技術名詞，建立高管與技術團隊的共同語言。"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "破冰：AI 在過去幾年如何改變世界？"
    tf = slide.placeholders[1].text_frame
    tf.text = "從自動駕駛到 ChatGPT：隱形 AI 到顯性 AI"
    tf.add_paragraph().text = "為何企業必須現在關注？「AI 不會取代你，但懂 AI 的人會」"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "技術名詞白話文翻譯"
    tf = slide.placeholders[1].text_frame
    tf.text = "人工智慧 (AI)：讓機器模仿人類智慧（大概念）"
    tf.add_paragraph().text = "機器學習 (ML)：給機器資料，讓它自己找規律（預測模型）"
    tf.add_paragraph().text = "深度學習 (DL)：用類神經網路處理複雜問題（影像、語音）"
    tf.add_paragraph().text = "生成式 AI (GenAI)：不僅能預測，還能「創造」新內容（GPT）"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "為什麼是現在？AI 的「奇點時刻」"
    tf = slide.placeholders[1].text_frame
    tf.text = "算力 (Computing Power)：GPU 技術的爆發（NVIDIA 等）"
    tf.add_paragraph().text = "數據 (Data)：網際網路累積的巨量訓練資料"
    tf.add_paragraph().text = "演算法 (Algorithm)：Transformer 模型架構的突破"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI 的能力邊界：能做什麼？不能做什麼？"
    tf = slide.placeholders[1].text_frame
    tf.text = "✅ 擅長：大量資料歸納、草稿生成、規律預測、跨語言翻譯"
    tf.add_paragraph().text = "❌ 不擅長：情緒共鳴、絕對的事實精準（幻覺問題）、承擔法律與道德責任"

    # --- Module 2 ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "【Module 2】2026 最熱門的 AI 話題與趨勢"
    slide.placeholders[1].text = "核心目標：掌握最新關鍵字，看懂科技巨頭的佈局方向。"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "大型語言模型 (LLM) 的進化與分化"
    tf = slide.placeholders[1].text_frame
    tf.text = "通用大模型：GPT-4、Gemini 3 等閉源巨獸"
    tf.add_paragraph().text = "開源模型崛起：Meta Llama 3 帶來的開源紅利"
    tf.add_paragraph().text = "企業私有模型：注重資安，用企業內部數據微調 (Fine-tuning) 的小模型"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "多模態 AI (Multimodal AI) 的全面爆發"
    tf = slide.placeholders[1].text_frame
    tf.text = "文字、圖片、聲音、影片的融合處理"
    tf.add_paragraph().text = "商業應用：從輸入一張設計圖直接生成程式碼，或用語音直接指揮系統操作"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "從 Copilot 到 AI Agents (自主代理)"
    tf = slide.placeholders[1].text_frame
    tf.text = "過去 (Copilot)：你下一個指令，它做一個動作（副駕駛）"
    tf.add_paragraph().text = "現在與未來 (Agents)：給予目標，AI 自己規劃步驟、上網找資料、執行工具（數位員工）"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "邊緣 AI 與 AI PC/手機"
    tf = slide.placeholders[1].text_frame
    tf.text = "算力下放：不需要連網也能執行 AI 模型"
    tf.add_paragraph().text = "優勢：解決企業最擔心的「機密資料外洩上雲」問題，大幅降低延遲"

    # --- Module 3 ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "【Module 3】企業實戰 —— AI 的商業價值與案例"
    slide.placeholders[1].text = "核心目標：探討 ROI，看看別人怎麼用 AI 賺錢或省錢。"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI 帶來的企業價值維度"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. 降本 (Cost Reduction)：自動化重複性勞動"
    tf.add_paragraph().text = "2. 增效 (Efficiency)：賦能員工，提升單兵作戰能力"
    tf.add_paragraph().text = "3. 創新 (Innovation)：創造過去無法實現的新產品或服務體驗"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "實戰案例 1：行銷與超級個人化體驗"
    tf = slide.placeholders[1].text_frame
    tf.text = "動態生成廣告素材：根據不同客群瞬間生成上千種文案與圖片"
    tf.add_paragraph().text = "智能客服 2.0：不再是死板的機器人，能理解情緒、查閱顧客歷史並安撫解決問題"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "實戰案例 2：打造企業專屬「知識大腦」"
    tf = slide.placeholders[1].text_frame
    tf.text = "痛點：員工離職帶走知識，新人培訓耗時"
    tf.add_paragraph().text = "解法：將公司過去 10 年的 SOP、合約、報價單餵給內部 AI"
    tf.add_paragraph().text = "效益：新人只需向 AI 提問「我們對 A 客戶的退貨規定是什麼？」，秒速解答"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "實戰案例 3：營運、製造與供應鏈"
    tf = slide.placeholders[1].text_frame
    tf.text = "需求預測：結合天氣、社群風向、歷史銷量預測下個月庫存需求"
    tf.add_paragraph().text = "機器視覺：產線上的瑕疵檢測，精準度與速度超越人工"
    tf.add_paragraph().text = "自動排程：針對複雜訂單動態調整生產線"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "實戰案例 4：軟體開發與跨部門數據分析"
    tf = slide.placeholders[1].text_frame
    tf.text = "工程師增效：GitHub Copilot 協助撰寫 40% 的常規程式碼"
    tf.add_paragraph().text = "全民數據分析師：業務主管不用會寫 SQL，只需用自然語言說「幫我畫出上季各區業績比較圖」，AI 自動撈資料出圖"

    # --- Module 4 ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "【Module 4】高管必修 —— AI 導入策略與風險管理"
    slide.placeholders[1].text = "核心目標：給予具體的行動指南 (Action Plan) 與避坑守則。"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "企業導入 AI 的四大陷阱"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. FOMO 焦慮：為了做而做，沒有鎖定真正的業務痛點"
    tf.add_paragraph().text = "2. 數據基建薄弱：垃圾進，垃圾出 (Garbage in, garbage out)"
    tf.add_paragraph().text = "3. 忽略員工抗拒：員工怕被取代而拒絕使用新工具"
    tf.add_paragraph().text = "4. 資安與合規風險：員工將機密財報丟上公開版 ChatGPT 翻譯"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "數據治理與資安防護策略"
    tf = slide.placeholders[1].text_frame
    tf.text = "建立企業內部 AI 使用守則 (AI Policy)"
    tf.add_paragraph().text = "採購企業版 AI (如 Copilot for Microsoft 365, ChatGPT Enterprise)，確保資料不被用於模型訓練"
    tf.add_paragraph().text = "資料分級制度：機密資料不上公有雲"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "組織變革：如何讓員工擁抱 AI？"
    tf = slide.placeholders[1].text_frame
    tf.text = "從 Top-down 支持，加上 Bottom-up 創新"
    tf.add_paragraph().text = "舉辦企業內部「AI 黑客松」或「最佳提示詞 (Prompt) 競賽」"
    tf.add_paragraph().text = "重新定義 KPI：獎勵那些利用 AI 省下時間並創造更高價值的員工"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "給各位的行動指南：導入 AI 的 100 天計畫"
    tf = slide.placeholders[1].text_frame
    tf.text = "Day 1-30：盤點痛點，挑選一個「低風險、高回報」的痛點作為起點 (POC)"
    tf.add_paragraph().text = "Day 31-60：選擇適合的工具與夥伴，小規模試點測試"
    tf.add_paragraph().text = "Day 61-90：收集回饋，修正流程，計算初步 ROI"
    tf.add_paragraph().text = "Day 100+：建立成功案例，推廣至全公司"

    # End
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Q&A 交流時間"
    slide.placeholders[1].text = "謝謝聆聽！歡迎提問探討。"

    prs.save("EMBA_AI_Strategy_Presentation.pptx")

if __name__ == "__main__":
    create_emba_presentation()
    print("EMBA Presentation PPTX Created.")
